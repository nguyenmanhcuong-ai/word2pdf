# app.py
from flask import Flask, render_template, request, send_file, redirect, url_for, jsonify
import os
from werkzeug.utils import secure_filename
import tempfile
import shutil
import pythoncom
import uuid
import time
from pdf2docx import Converter as PDFToDocxConverter
from docx2pdf import convert as docx_to_pdf_convert
import comtypes.client
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import zipfile
import logging

# Thiết lập logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024  # 32MB max file size
app.config['CLEANUP_THRESHOLD'] = 3600  # 1 giờ (tính bằng giây)

# Đảm bảo thư mục uploads và outputs tồn tại
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

ALLOWED_EXTENSIONS = {
    'pdf': ['docx', 'pptx', 'txt', 'jpg', 'png'],
    'docx': ['pdf'],
    'pptx': ['pdf'],
    'jpg': ['pdf'],
    'jpeg': ['pdf'],
    'png': ['pdf'],
    'txt': ['pdf']
}

def allowed_file(filename):
    """Kiểm tra tệp có đúng định dạng được hỗ trợ không"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS.keys()

def get_possible_conversions(file_ext):
    """Lấy danh sách các định dạng mà file_ext có thể chuyển đổi sang"""
    file_ext = file_ext.lower()
    if file_ext in ALLOWED_EXTENSIONS:
        return ALLOWED_EXTENSIONS[file_ext]
    return []

@app.route('/')
def index():
    return render_template('index.html', conversion_types=ALLOWED_EXTENSIONS)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'Không có tệp nào được chọn'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'Không có tệp nào được chọn'}), 400
    
    if file and allowed_file(file.filename):
        # Tạo id cho tệp để tránh xung đột tên
        file_id = str(uuid.uuid4())
        filename = secure_filename(file.filename)
        file_ext = filename.rsplit('.', 1)[1].lower()
        
        # Lưu tệp với id
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        file.save(file_path)
        
        possible_conversions = get_possible_conversions(file_ext)
        return jsonify({
            'success': True,
            'file_id': file_id,
            'filename': filename,
            'file_ext': file_ext,
            'possible_conversions': possible_conversions
        })
    
    return jsonify({'error': 'Định dạng tệp không được hỗ trợ'}), 400

@app.route('/convert', methods=['POST'])
def convert_file():
    try:
        data = request.get_json()
        file_id = data.get('file_id')
        filename = data.get('filename')
        target_format = data.get('target_format')
        
        if not all([file_id, filename, target_format]):
            return jsonify({'error': 'Thiếu thông tin cần thiết'}), 400
        
        orig_file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        if not os.path.exists(orig_file_path):
            return jsonify({'error': 'Tệp không tồn tại'}), 404
        
        file_ext = filename.rsplit('.', 1)[1].lower()
        output_filename = f"{filename.rsplit('.', 1)[0]}.{target_format}"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{output_filename}")
        
        # Gọi hàm chuyển đổi tương ứng
        if file_ext == 'pdf' and target_format == 'docx':
            pdf_to_docx(orig_file_path, output_path)
        elif file_ext == 'docx' and target_format == 'pdf':
            docx_to_pdf(orig_file_path, output_path)
        elif file_ext == 'pptx' and target_format == 'pdf':
            pptx_to_pdf(orig_file_path, output_path)
        elif file_ext == 'pdf' and target_format == 'pptx':
            pdf_to_pptx(orig_file_path, output_path)
        elif file_ext in ['jpg', 'jpeg', 'png'] and target_format == 'pdf':
            image_to_pdf(orig_file_path, output_path)
        elif file_ext == 'pdf' and target_format in ['jpg', 'png']:
            pdf_to_images(orig_file_path, output_path, target_format)
        elif file_ext == 'txt' and target_format == 'pdf':
            txt_to_pdf(orig_file_path, output_path)
        elif file_ext == 'pdf' and target_format == 'txt':
            pdf_to_txt(orig_file_path, output_path)
        else:
            return jsonify({'error': 'Loại chuyển đổi không được hỗ trợ'}), 400
        
        download_url = url_for('download_file', file_id=file_id, filename=output_filename)
        return jsonify({'success': True, 'download_url': download_url})
        
    except Exception as e:
        logger.error(f"Lỗi khi chuyển đổi: {str(e)}")
        return jsonify({'error': f"Lỗi khi chuyển đổi: {str(e)}"}), 500

@app.route('/download/<file_id>/<filename>')
def download_file(file_id, filename):
    file_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{filename}")
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True, download_name=filename)
    return "Tệp không tồn tại", 404

# Các hàm chuyển đổi 
def pdf_to_docx(pdf_path, docx_path):
    """Chuyển đổi PDF sang DOCX"""
    cv = PDFToDocxConverter(pdf_path)
    cv.convert(docx_path)
    cv.close()
    return docx_path

def docx_to_pdf(docx_path, pdf_path):
    """Chuyển đổi DOCX sang PDF"""
    pythoncom.CoInitialize()  # Cần thiết cho Windows COM
    docx_to_pdf_convert(docx_path, pdf_path)
    return pdf_path

def pptx_to_pdf(pptx_path, pdf_path):
    """Chuyển đổi PPTX sang PDF"""
    pythoncom.CoInitialize()
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = True
    
    try:
        deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
        deck.SaveAs(os.path.abspath(pdf_path), 32)  # 32 là mã định dạng PDF
        deck.Close()
    finally:
        powerpoint.Quit()
    
    return pdf_path

def pdf_to_pptx(pdf_path, pptx_path):
    """Chuyển đổi PDF sang PPTX"""
    # Mở PDF và chuyển từng trang thành ảnh
    pdf_document = fitz.open(pdf_path)
    prs = Presentation()
    
    for page_number in range(pdf_document.page_count):
        page = pdf_document[page_number]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Tăng độ phân giải
        
        # Lưu ảnh tạm thời
        temp_img_path = os.path.join(tempfile.gettempdir(), f"page_{page_number}.png")
        pix.save(temp_img_path)
        
        # Thêm slide mới và chèn ảnh
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slide.shapes.add_picture(temp_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Xóa tệp tạm
        os.remove(temp_img_path)
    
    prs.save(pptx_path)
    pdf_document.close()
    
    return pptx_path

def image_to_pdf(image_path, pdf_path):
    """Chuyển đổi ảnh sang PDF"""
    image = Image.open(image_path)
    # Chuyển sang RGB nếu là RGBA (để tránh lỗi khi lưu PDF)
    if image.mode == 'RGBA':
        image = image.convert('RGB')
    image.save(pdf_path, 'PDF')
    return pdf_path

def pdf_to_images(pdf_path, output_path, format='jpg'):
    """Chuyển đổi PDF sang các tệp ảnh"""
    # Thay đổi đuôi file thành zip vì sẽ có nhiều ảnh
    zip_path = output_path.rsplit('.', 1)[0] + '.zip'
    
    # Mở PDF
    pdf_document = fitz.open(pdf_path)
    
    # Tạo thư mục tạm để chứa ảnh
    temp_dir = tempfile.mkdtemp()
    
    try:
        with zipfile.ZipFile(zip_path, 'w') as zip_file:
            for page_number in range(pdf_document.page_count):
                page = pdf_document[page_number]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Tăng độ phân giải
                
                # Lưu ảnh vào thư mục tạm
                img_filename = f"page_{page_number+1}.{format}"
                img_path = os.path.join(temp_dir, img_filename)
                pix.save(img_path)
                
                # Thêm ảnh vào file zip
                zip_file.write(img_path, img_filename)
    
    finally:
        # Dọn dẹp
        pdf_document.close()
        shutil.rmtree(temp_dir)
    
    return zip_path

def txt_to_pdf(txt_path, pdf_path):
    """Chuyển đổi TXT sang PDF"""
    # Đọc nội dung tệp văn bản
    with open(txt_path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    # Tạo PDF từ đoạn văn bản
    doc = fitz.open()
    page = doc.new_page()
    
    # Thêm văn bản vào trang
    page.insert_text((50, 50), text, fontsize=11)
    
    # Lưu PDF
    doc.save(pdf_path)
    doc.close()
    
    return pdf_path

def pdf_to_txt(pdf_path, txt_path):
    """Trích xuất văn bản từ PDF"""
    # Mở PDF
    doc = fitz.open(pdf_path)
    
    # Trích xuất văn bản từ tất cả các trang
    text = ""
    for page_num in range(doc.page_count):
        page = doc[page_num]
        text += page.get_text()
    
    # Lưu văn bản vào tệp
    with open(txt_path, 'w', encoding='utf-8') as f:
        f.write(text)
    
    doc.close()
    return txt_path

@app.route('/cleanup', methods=['POST'])
def cleanup_files():
    """Dọn dẹp tệp cũ"""
    current_time = time.time()
    count = 0
    
    # Dọn thư mục uploads
    for filename in os.listdir(app.config['UPLOAD_FOLDER']):
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        if os.path.isfile(file_path):
            # Kiểm tra thời gian tạo tệp
            file_time = os.path.getmtime(file_path)
            if current_time - file_time > app.config['CLEANUP_THRESHOLD']:
                os.remove(file_path)
                count += 1
    
    # Dọn thư mục outputs
    for filename in os.listdir(app.config['OUTPUT_FOLDER']):
        file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        if os.path.isfile(file_path):
            file_time = os.path.getmtime(file_path)
            if current_time - file_time > app.config['CLEANUP_THRESHOLD']:
                os.remove(file_path)
                count += 1
    
    return jsonify({'success': True, 'deleted_files': count})

if __name__ == '__main__':
    app.run(debug=True)